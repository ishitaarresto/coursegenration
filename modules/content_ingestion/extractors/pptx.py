"""
PPTX extractor using python-pptx.

Responsibilities:
- Open the presentation via pptx.Presentation(path)
- Iterate slides; for each slide sort all TEXT_BOX / text-bearing shapes by
  (top, left) to preserve visual reading order
- Extract speaker notes from slide.notes_slide.notes_text_frame
- Extract embedded images via the slide's XML relationship table rather than
  by shape type -- images in PPTX can be stored as standalone PICTURE shapes
  (MSO_SHAPE_TYPE.PICTURE), as picture fills on FREEFORM / AUTO_SHAPE shapes,
  or inside GROUP shapes; all of these reference their image blobs through the
  slide part's relationship table (rId -> image part), so iterating rels is the
  only approach that catches all of them
- Deduplicate images by content fingerprint so the same texture used on
  multiple shapes is captioned only once per slide
- Populate ExtractedContent.slides with one ExtractedSlide per slide

Reading order: python-pptx returns shapes in XML document order (creation
order). Sorting by (top, left) approximates visual top-to-bottom reading
order for typical slide layouts.
"""

import io

from pptx import Presentation

from modules.content_ingestion.extractors.base import BaseExtractor
from modules.content_ingestion.models import (
    Asset, AssetType, ExtractedContent, ExtractedImage, ExtractedSlide,
)

MIN_IMAGE_DIM = 100  # pixels; skip images smaller than this (icons, bullets)


def _extract_slide_images(slide, start_index: int) -> list[ExtractedImage]:
    """Extract all images from a slide via its part relationship table.

    Every image a slide references -- whether it's a standalone picture shape,
    a picture fill on a freeform shape, or a photo inside a group shape -- is
    registered as an image relationship on the slide's XML part. Iterating
    those relationships is the only reliable way to catch all of them.

    Deduplicates by hashing the first 1 KB of each blob so the same background
    texture used on ten shapes is only extracted once.
    """
    from PIL import Image as PILImage

    images: list[ExtractedImage] = []
    seen_fingerprints: set[int] = set()
    idx = start_index

    for rel in slide.part.rels.values():
        if "image" not in rel.reltype.lower():
            continue
        try:
            blob: bytes = rel.target_part.blob
            mime: str = rel.target_part.content_type or ""

            # Deduplicate identical blobs (same texture referenced multiple times)
            fingerprint = hash(blob[:1024])
            if fingerprint in seen_fingerprints:
                continue
            seen_fingerprints.add(fingerprint)

            pil_img = PILImage.open(io.BytesIO(blob))
            w, h = pil_img.size

            if w < MIN_IMAGE_DIM or h < MIN_IMAGE_DIM:
                continue

            images.append(ExtractedImage(
                index=idx,
                image_bytes=blob,
                width=w,
                height=h,
                mime_type=mime,
            ))
            idx += 1
        except Exception:
            continue

    return images


class PptxExtractor(BaseExtractor):
    """Extracts per-slide text, speaker notes, and embedded images from PPTX."""

    def __init__(self, extract_images: bool = True) -> None:
        self.extract_images = extract_images

    def can_handle(self, asset: Asset) -> bool:
        return asset.asset_type == AssetType.PPTX

    def extract(self, asset: Asset) -> ExtractedContent:
        """Iterate slides; collect shape text in reading order plus images."""
        self._validate_file(asset)
        result = ExtractedContent(asset=asset)

        try:
            prs = Presentation(asset.file_path)
        except Exception as exc:
            result.extraction_errors.append(f"Presentation open failed: {exc}")
            return result

        slides_out: list[ExtractedSlide] = []

        for i, slide in enumerate(prs.slides):
            # -- Text: sort text-bearing shapes by visual position ----------
            text_shapes = sorted(
                (s for s in slide.shapes if s.has_text_frame),
                key=lambda s: (
                    s.top  if s.top  is not None else 0,
                    s.left if s.left is not None else 0,
                ),
            )
            lines: list[str] = []
            for shape in text_shapes:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)

            # -- Images: extract all via slide relationships -----------------
            slide_images: list[ExtractedImage] = []
            if self.extract_images:
                slide_images = _extract_slide_images(slide, start_index=0)
                # Append [Image N] placeholders after the slide text.
                # Images are fills / group members without fixed text-stream
                # positions, so they go at the end of the slide block.
                for img in slide_images:
                    lines.append(f"[Image {img.index}]")

            raw_text = "\n".join(lines)

            # -- Speaker notes ----------------------------------------------
            notes_text = ""
            try:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            slides_out.append(ExtractedSlide(
                slide_number=i + 1,
                raw_text=raw_text,
                speaker_notes=notes_text,
                images=slide_images,
            ))

        result.slides = slides_out
        result.full_text = "\n\n".join(s.raw_text for s in slides_out if s.raw_text)
        return result
